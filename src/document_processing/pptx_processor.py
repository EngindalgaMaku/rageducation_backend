"""
PPTX document processing module.

This module provides functions to extract text from PPTX files using python-pptx.
"""

from pptx import Presentation
from ..utils.helpers import setup_logging

logger = setup_logging()

def process_pptx(file_path: str) -> str:
    """
    Extracts text from a PPTX file.

    Args:
        file_path: The path to the PPTX file.

    Returns:
        The extracted text as a single string, or an empty string if an error occurs.
    """
    logger.info(f"Processing PPTX file: {file_path}")
    text = ""
    try:
        presentation = Presentation(file_path)
        for i, slide in enumerate(presentation.slides):
            # Inject a slide marker to preserve structure
            text += f"\n\n=== Slide {i+1} ===\n\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        logger.info(f"Successfully extracted text from {file_path}")
        return text.strip()
    except FileNotFoundError:
        logger.error(f"File not found: {file_path}")
        return ""
    except Exception as e:
        logger.error(f"An error occurred while processing {file_path}: {e}")
        return ""

if __name__ == '__main__':
    # This is for testing purposes.
    # Create a dummy PPTX file to test the processor.
    import os
    if not os.path.exists('data/raw/sample_documents'):
        os.makedirs('data/raw/sample_documents')

    test_pptx_path = "data/raw/sample_documents/sample.pptx"
    
    # Create a dummy pptx file for testing
    prs = Presentation()
    slide_layout = prs.slide_layouts  # A layout with a title and content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Test Presentation"
    
    content_shape = slide.placeholders
    content_shape.text = "This is a test slide in a PPTX document for the RAG system."
    
    prs.save(test_pptx_path)

    if os.path.exists(test_pptx_path):
        extracted_text = process_pptx(test_pptx_path)
        if extracted_text:
            print("--- Extracted Text from PPTX ---")
            print(extracted_text)
            print("------------------------------")
        else:
            print("Could not extract text from the PPTX.")
    else:
        print(f"Test file not found: {test_pptx_path}. Please create it to run the test.")
